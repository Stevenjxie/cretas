"""
用 Playwright 截取 slides.html 中每张幻灯片 → PNG
然后用 python-pptx 组装成 PPTX
"""
import asyncio, os, pathlib, sys
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches, Emu

BASE  = pathlib.Path(__file__).parent
HTML  = (BASE / "slides.html").as_uri()  # file:///...
OUTD  = BASE / "assets" / "slide_pngs"
PPTX  = BASE / "haijuying-cretas-pitch.pptx"

SLIDE_IDS = ["s0", "s1", "s2", "s3", "s4", "s4b", "s4c", "s4d", "s4e", "s5", "s6", "s7", "send"]
SLIDE_W = 1280
SLIDE_H = 720

OUTD.mkdir(parents=True, exist_ok=True)


async def screenshot_all():
    async with async_playwright() as p:
        browser = await p.chromium.launch(
            headless=False,
            args=[
                "--lang=zh-CN",
                "--font-render-hinting=none",
                "--disable-blink-features=AutomationControlled",
            ]
        )
        page = await browser.new_page(viewport={"width": SLIDE_W, "height": SLIDE_H})
        await page.goto(HTML, wait_until="networkidle")
        await page.wait_for_timeout(800)  # let fonts settle

        paths = []
        for slide_id in SLIDE_IDS:
            el = await page.query_selector(f"#{slide_id}")
            if el is None:
                print(f"⚠️  #{slide_id} not found, skipping")
                continue
            # scroll the element into view and screenshot just this element
            await el.scroll_into_view_if_needed()
            out = OUTD / f"{slide_id}.png"
            await el.screenshot(path=str(out))
            print(f"OK  {slide_id} -> {out.name}")
            paths.append((slide_id, str(out)))

        await browser.close()
        return paths


def build_pptx(png_paths):
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]  # blank

    for slide_id, png_path in png_paths:
        sl = prs.slides.add_slide(blank)
        # fill the slide background white (won't show — image covers it)
        sl.shapes.add_picture(
            png_path,
            left=0, top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(str(PPTX))
    size_kb = PPTX.stat().st_size // 1024
    print(f"\nPPTX saved: {PPTX}")
    print(f"   Slides: {len(prs.slides)}  /  Size: {size_kb} KB")


async def main():
    print("Screenshotting slides...")
    paths = await screenshot_all()

    print("\nBuilding PPTX...")
    build_pptx(paths)


if __name__ == "__main__":
    asyncio.run(main())
